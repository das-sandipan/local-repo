from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()

# Define slide layout
title_slide_layout = prs.slide_layouts[5]  # Title + Content

# Helper function to add slide
def add_slide(title, bullets, code=None):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[0]

    # Title
    title_shape.text = title

    # Bullets
    tf = body_shape.text_frame
    tf.clear()
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)

    # Code box
    if code:
        left = Inches(5)
        top = Inches(2)
        width = Inches(4.5)
        height = Inches(3)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf_code = textbox.text_frame
        tf_code.text = code
        for p in tf_code.paragraphs:
            p.font.name = "Courier New"
            p.font.size = Pt(16)

# Slides
add_slide("Introduction", [
    "Python is an open-source general purpose scripting language.",
    "Written in C Language.",
    "Used for websites (Django, Flask) and machine learning (Scikit-learn).",
    "Interpreted, not compiled.",
    "Developed by Guido van Rossum, first released in 1991."
])

add_slide("Python Versions", [
    "Python 3.x widely used; Python 2.x mostly forward compatible.",
    "Packages for data analysis available in both versions.",
    "Division difference: Python 2 (5/2=2), Python 3 (5/2=2.5).",
    "Multiple versions/environments can co-exist (Conda, venv).",
    "Anaconda bundles useful packages and IDEs like JupyterLab."
])

add_slide("Basic Syntax", [
    "Variables can be printed directly or as strings.",
    "Use help(x) to inspect objects.",
    "Indentation defines code blocks.",
    "Line breaks use backslash \\."
], code="""x = 5
print("x")   # prints literal string
print(x)     # prints variable value""")

add_slide("Strings", [
    "Strings are characters within quotes.",
    "Indexed forward (0,1,2,…) and backward (-1,-2,…).",
    "Operators: + concatenation, * repetition, %s formatting."
], code="""str = "Hello"
str[0]      # 'H'
str[-1]     # 'o'
str[1:3]    # 'el'
print("%s cannot be %s" %("the moon", "stolen"))""")

add_slide("Lists", [
    "Defined with square brackets [ ].",
    "Can contain mixed items.",
    "Indexed, iterable, mutable.",
    "Operators: + concatenation, * repetition.",
    "Functions: append, insert, extend, pop."
], code="""myList = ["Movie", "Mathematics", "Magic"]
myList[0]           # 'Movie'
myList[:2]          # ['Movie', 'Mathematics']
myList[-1]          # 'Magic'
myList + ["Biking"] # concatenation
myList * 2          # repetition""")

# Save presentation
prs.save("Python_Basics.pptx")
print("Presentation saved as Python_Basics.pptx")